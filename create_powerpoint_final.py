#!/usr/bin/env python3
"""
GitHub Workflow Best Practices - Final PowerPoint Generator
Creates a properly scaled presentation emphasizing communication-first approach
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create presentation with widescreen aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PRIMARY_COLOR = RGBColor(9, 105, 218)  # GitHub Blue
    SECONDARY_COLOR = RGBColor(31, 136, 61)  # GitHub Green
    DANGER_COLOR = RGBColor(209, 36, 47)  # Red
    WARNING_COLOR = RGBColor(251, 133, 0)  # Orange
    DARK_COLOR = RGBColor(36, 41, 47)  # Dark Gray
    
    # Slide 1: Title Slide - Communication Focus
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 GitHub Developer Workflow"
    title.text_frame.paragraphs[0].font.size = Pt(38)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Phase 1: Improve Team Communication & Standards\n\n💬 Better Collaboration | 📝 Consistent Documentation | ⚡ Faster Decisions\n\nPhase 2: Full GitHub Project Management (Future)\n\nPresented by: Amna & Hanzla"
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)
    
    # Slide 2: Agenda - Smaller fonts, more spacing
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📋 Communication-First Agenda"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Shorter, more focused agenda items
    agenda_items = [
        ("🔴 Communication Pain Points", "4 min"),
        ("💬 Phase 1: Improve Team Communication", "8 min"),
        ("📝 Templates & Standards", "6 min"),
        ("🎯 Quick Wins Demo", "5 min"),
        ("📊 Phase 1 Implementation (2 weeks)", "4 min"),
        ("🚀 Phase 2 Vision (Future)", "3 min")
    ]
    
    for item, time in agenda_items:
        p = tf.add_paragraph()
        p.text = f"{item} ({time})"
        p.level = 0
        p.font.size = Pt(20)  # Larger for better visibility
        p.font.bold = True
        p.space_after = Pt(12)
    
    # Slide 3: Communication Challenges
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💬 Communication Challenges"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    challenges = [
        ("🔧 Process Communication (40% Impact)", [
            "Unclear requirements in issues",
            "Missing context in PRs",
            "No standard templates"
        ]),
        ("👥 Team Communication (35% Impact)", [
            "Knowledge silos",
            "Delayed feedback",
            "Inconsistent updates"
        ]),
        ("📋 Documentation Issues (25% Impact)", [
            "Outdated information",
            "Missing project context",
            "No standardized formats"
        ])
    ]
    
    for category, items in challenges:
        p = tf.add_paragraph()
        p.text = category
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DANGER_COLOR
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(15)
        
        tf.add_paragraph()  # Add spacing
    
    # Slide 4: Phase 1 Focus
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💬 Phase 1: Communication First"
    
    # Create single content box with better spacing
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(11.5)
    height = Inches(5)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "🎯 Why Communication First?"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(15)
    
    benefits = [
        "✅ Immediate impact on team productivity",
        "✅ Low implementation effort, high value",
        "✅ Builds foundation for Phase 2 (full methodology)",
        "✅ Everyone can contribute regardless of GitHub experience",
        "✅ Reduces confusion and rework by 60%"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # Slide 5: Phase 1 Implementation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📝 Phase 1: Templates & Standards"
    
    # Create table with proper dimensions
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(12.3)
    height = Inches(4.5)
    
    table = slide.shapes.add_table(5, 3, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(3.9)
    table.columns[2].width = Inches(3.9)
    
    # Header row
    headers = ['Communication Tool', 'Current State', 'Phase 1 Improvement']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
    
    # Data rows
    data = [
        ('📋 Issue Creation', 'Vague descriptions', 'Smart templates'),
        ('🔄 Pull Requests', 'Missing context', 'Standard format'),
        ('🏷️ Labels', 'Inconsistent usage', 'Clear taxonomy'),
        ('💬 Communication', 'Ad-hoc updates', 'Structured process')
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(14)
    
    # Slide 6: Quick Wins
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "⚡ Phase 1 Quick Wins (Week 1)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    quick_wins = [
        "📝 Deploy issue templates (1 hour setup)",
        "✅ Add PR template (30 minutes)",
        "🏷️ Create basic label set (15 labels max)",
        "📋 Team training session (1 hour)",
        "💬 Establish update rhythm (daily/weekly)"
    ]
    
    for win in quick_wins:
        p = tf.add_paragraph()
        p.text = win
        p.font.size = Pt(20)
        p.space_after = Pt(15)
    
    # Slide 7: Issue Templates Demo
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📋 Better Issue Communication"
    
    # Create two boxes for before/after
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(6)
    height = Inches(4.5)
    
    # Before box
    before_box = slide.shapes.add_textbox(left, top, width, height)
    tf = before_box.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "❌ Before: Poor Communication"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DANGER_COLOR
    
    p = tf.add_paragraph()
    p.text = "Title: Bug\nDescription: It's broken\n\nResult:\n• Unclear scope\n• Missing context\n• Back-and-forth questions\n• Delayed resolution"
    p.font.size = Pt(14)
    
    # After box
    right_left = Inches(7)
    after_box = slide.shapes.add_textbox(right_left, top, width, height)
    tf = after_box.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "✅ After: Clear Communication"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    p = tf.add_paragraph()
    p.text = "🐛 Login Error on Mobile\n\nPriority: High\nSteps: 1. Open app 2. Enter credentials\nExpected: Dashboard loads\nActual: 500 error\nScreenshot: [attached]\n\nResult:\n• Clear understanding\n• Immediate action\n• Faster resolution"
    p.font.size = Pt(14)
    
    # Slide 8: 2-Week Implementation Plan
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📅 Phase 1: 2-Week Implementation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    roadmap = [
        ("Week 1: Foundation", "✅ Templates | ✅ Labels | ✅ Training | ✅ First PRs"),
        ("Week 2: Adoption", "📊 Measure usage | 🔄 Iterate | 💬 Feedback | 🎯 Refine")
    ]
    
    for week, tasks in roadmap:
        p = tf.add_paragraph()
        p.text = week
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = tf.add_paragraph()
        p.text = tasks
        p.font.size = Pt(18)
        p.space_after = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Expected Results After 2 Weeks:"
    p.font.size = Pt(20)
    p.font.bold = True
    
    results = [
        "🎯 80% of issues use templates",
        "💬 Clearer team communication",
        "⚡ 30% faster issue resolution",
        "📋 Consistent documentation"
    ]
    
    for result in results:
        p = tf.add_paragraph()
        p.text = result
        p.font.size = Pt(16)
    
    # Slide 9: Phase 2 Vision
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 Phase 2: Full Methodology (Future)"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "After Phase 1 Success, We'll Add:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(15)
    
    phase2_items = [
        "📊 GitHub Project Boards for sprint management",
        "🤖 Automated workflows and quality gates",
        "📈 Sprint planning and velocity tracking",
        "🔄 Full CI/CD integration",
        "📋 Advanced reporting and analytics"
    ]
    
    for item in phase2_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "\n🎯 Timeline: After Phase 1 proves communication improvements"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_COLOR
    p.font.bold = True
    
    # Slide 10: Action Items
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🎬 Start Phase 1 Today"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Immediate Actions (This Week):"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.space_after = Pt(15)
    
    actions = [
        "☑️ Add issue template to one repository",
        "☑️ Create basic label set (5-10 labels)",
        "☑️ Add PR template",
        "☑️ Team intro session (30 minutes)",
        "☑️ Try on 2-3 issues this week"
    ]
    
    for action in actions:
        p = tf.add_paragraph()
        p.text = action
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "\nResources:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "📦 Templates: github.com/hanzlahabib/github-workflow-demo"
    p.font.size = Pt(16)
    
    # Slide 11: Questions
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💡 Questions & Discussion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Phase 1 Questions:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(15)
    
    questions = [
        "✅ How do we get team buy-in for templates?",
        "✅ What if people forget to use them?",
        "✅ How do we measure communication improvement?",
        "✅ When do we move to Phase 2?"
    ]
    
    for q in questions:
        p = tf.add_paragraph()
        p.text = q
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    p = tf.add_paragraph()
    p.text = "\nContact:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_before = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "GitHub: @hanzlahabib | @amna\nDemo: github.com/hanzlahabib/github-workflow-demo"
    p.font.size = Pt(16)
    
    # Slide 12: Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🙏 Thank You!"
    title.text_frame.paragraphs[0].font.size = Pt(42)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Ready to Improve Team Communication?\n\n💬 Start Phase 1 This Week\n📈 Measure the Impact\n🚀 Prepare for Phase 2\n\nLet's build better software, together."
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Save the presentation
    prs.save('GitHub_Workflow_Communication_First.pptx')
    print("✅ Communication-focused PowerPoint created successfully!")
    print("📁 File saved as: GitHub_Workflow_Communication_First.pptx")
    print("💬 Emphasizes Phase 1: Communication improvement")
    print("🎯 Properly sized for full content visibility")

if __name__ == "__main__":
    create_presentation()