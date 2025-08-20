#!/usr/bin/env python3
"""
GitHub Workflow Best Practices - Improved PowerPoint Generator
Creates a professional PowerPoint presentation with optimized sizing and layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create presentation with widescreen aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color scheme
    PRIMARY_COLOR = RGBColor(9, 105, 218)  # GitHub Blue
    SECONDARY_COLOR = RGBColor(31, 136, 61)  # GitHub Green
    DANGER_COLOR = RGBColor(209, 36, 47)  # Red
    WARNING_COLOR = RGBColor(251, 133, 0)  # Orange
    DARK_COLOR = RGBColor(36, 41, 47)  # Dark Gray
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 GitHub Developer Workflow"
    title.text_frame.paragraphs[0].font.size = Pt(42)  # Reduced from 54
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    subtitle = slide.placeholders[1]
    subtitle.text = "From Chaos to Clarity: Two-Phase Transformation\n\n🎯 Phase 1: Improve Communication & Standards\n🚀 Phase 2: Full GitHub Project & Agile Methodology\n\n📈 60% Faster PR Cycles | 40% Fewer Bugs | 90% Team Satisfaction\n\nPresented by: Amna & Hanzla"
    subtitle.text_frame.paragraphs[0].font.size = Pt(16)  # Reduced further to accommodate more text
    
    # Slide 2: Agenda
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📋 Agenda"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    agenda_items = [
        ("🔴 Current Pain Points", "5 min", "Identify workflow challenges"),
        ("🟡 Root Cause Analysis", "3 min", "Understanding persistence"),
        ("🟢 Solution Architecture", "10 min", "GitHub-based workflow system"),
        ("🎯 Live Demonstration", "7 min", "See workflow in action"),
        ("📊 Implementation Roadmap", "5 min", "30-day transformation plan")
    ]
    
    for item, time, desc in agenda_items:
        p = tf.add_paragraph()
        p.text = f"{item} ({time})"
        p.level = 0
        p.font.size = Pt(17)  # Reduced from 20
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)  # Reduced from 16
        p.space_after = Pt(8)  # Reduced from 12
    
    # Slide 3: Challenges
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚨 Current Challenges"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    challenges = [
        ("🔧 Process Issues (40% Impact)", [
            "Split Project Management - Multiple tools",
            "Inconsistent Formats - 30+ min/PR overhead",
            "Unclear Sprint Boundaries - 25% scope creep"
        ]),
        ("👥 People Issues (35% Impact)", [
            "Knowledge Gap - Single points of failure",
            "Poor Communication - Delayed decisions",
            "Review Quality vs Delay - 3-day PR lifecycle"
        ]),
        ("🛠️ Technical Debt (25% Impact)", [
            "Production Bugs - Missing quality gates",
            "Outdated Docs - 60% onboarding confusion",
            "Notification Chaos - Important items missed"
        ])
    ]
    
    for category, items in challenges:
        p = tf.add_paragraph()
        p.text = category
        p.level = 0
        p.font.size = Pt(16)  # Reduced from 20
        p.font.bold = True
        p.font.color.rgb = DANGER_COLOR
        
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(13)  # Reduced from 16
        
        tf.add_paragraph()  # Add spacing
    
    # Slide 4: Transformation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🔄 From Chaos to Clarity"
    
    # Create two text boxes with improved dimensions
    left = Inches(0.3)
    top = Inches(1.6)
    width = Inches(6.4)  # Increased width
    height = Inches(5.2)  # Increased height
    
    # Before box
    left_content = slide.shapes.add_textbox(left, top, width, height)
    tf = left_content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "😰 Before (Chaos)"
    p.font.size = Pt(18)  # Reduced from 24
    p.font.bold = True
    p.font.color.rgb = DANGER_COLOR
    
    before_items = [
        "5+ tools for project management",
        "3-day PR review cycle",
        "40% PRs missing context",
        "Weekly production bugs",
        "60% sprint completion rate"
    ]
    
    for item in before_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(15)  # Reduced from 18
        p.space_after = Pt(4)  # Reduced from 6
    
    # After box
    right_left = Inches(6.9)  # Adjusted position
    right_content = slide.shapes.add_textbox(right_left, top, width, height)
    tf = right_content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "🎯 After (Clarity)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    after_items = [
        "Single GitHub workspace",
        "4-hour PR review cycle",
        "100% PRs with templates",
        "90% bug reduction",
        "92% sprint completion rate"
    ]
    
    for item in after_items:
        p = tf.add_paragraph()
        p.text = f"✅ {item}"
        p.font.size = Pt(15)
        p.space_after = Pt(4)
    
    # Slide 5: Solution Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💡 Solution Architecture"
    
    # Create table with improved dimensions
    left = Inches(0.3)
    top = Inches(1.6)
    width = Inches(12.7)  # Increased width
    height = Inches(4.8)  # Increased height
    
    table = slide.shapes.add_table(6, 4, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(4.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(2.9)
    
    # Header row
    headers = ['Component', 'Impact', 'Effort', 'Timeline']
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(15)  # Reduced from 18
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
    
    # Data rows
    data = [
        ('📊 GitHub Project Boards', '⭐⭐⭐ High', 'Low', 'Week 1'),
        ('📝 Issue Templates', '⭐⭐⭐ High', 'Low', 'Week 1'),
        ('✅ PR Templates', '⭐⭐⭐ High', 'Low', 'Week 1'),
        ('🤖 GitHub Actions', '⭐⭐⭐ High', 'Medium', 'Week 2'),
        ('🏷️ Label Taxonomy', '⭐⭐ Medium', 'Low', 'Week 1')
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(13)  # Reduced from 16
    
    # Slide 6: Workflow Automation
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "⚙️ Automated Workflow"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    workflow_steps = [
        "📝 Issue Created → Template enforced",
        "🏷️ Auto-labeled → Type & Priority assigned",
        "📊 Added to Board → Sprint assigned automatically",
        "🌿 Branch Created → Follows naming convention",
        "🔄 PR Opened → Template applied, checks triggered",
        "✅ Auto-checks → Quality gates enforced",
        "🚀 Merged → Automatic deployment"
    ]
    
    for step in workflow_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(16)  # Reduced from 20
        p.space_after = Pt(8)  # Reduced from 12
    
    # Slide 7: Label System
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🏷️ Intelligent Label Taxonomy"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Type Labels (One Required):"
    p.font.size = Pt(17)  # Reduced from 20
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🟦 type: feature   🟥 type: bug   🟨 type: task"
    p.font.size = Pt(15)  # Reduced from 18
    p.space_after = Pt(15)  # Reduced from 20
    
    p = tf.add_paragraph()
    p.text = "Priority Matrix:"
    p.font.size = Pt(17)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🔴 critical   🟠 high   🟡 medium   🟢 low"
    p.font.size = Pt(15)
    p.space_after = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "Workflow Status (Auto-updated):"
    p.font.size = Pt(17)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "📋 backlog → 📝 todo → 🏃 in-progress"
    p.font.size = Pt(15)
    
    p = tf.add_paragraph()
    p.text = "🚫 blocked → 👀 review → ✅ done"
    p.font.size = Pt(15)
    
    # Slide 8: Implementation Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🚀 30-Day Implementation Plan"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    roadmap = [
        ("Week 1: Foundation", "✅ Deploy templates | ✅ Configure labels | ✅ Team training"),
        ("Week 2: Automation", "⚙️ GitHub Actions | 📊 Project boards | 🔗 Integrations"),
        ("Week 3: Integration", "💬 Slack notifications | 📈 Analytics dashboard"),
        ("Week 4: Optimization", "🎯 Refine workflows | 📊 Measure success | 🚀 Scale")
    ]
    
    for week, tasks in roadmap:
        p = tf.add_paragraph()
        p.text = week
        p.font.size = Pt(18)  # Reduced from 22
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = tf.add_paragraph()
        p.text = tasks
        p.font.size = Pt(15)  # Reduced from 18
        p.space_after = Pt(12)  # Reduced from 20
    
    # Slide 9: Success Metrics
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📈 Measurable Success"
    
    # Before metrics box with improved dimensions
    left = Inches(0.3)
    top = Inches(1.6)
    width = Inches(6.4)
    height = Inches(5.2)
    
    left_content = slide.shapes.add_textbox(left, top, width, height)
    tf = left_content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "📊 Before"
    p.font.size = Pt(20)  # Reduced from 24
    p.font.bold = True
    
    metrics_before = [
        "PR Cycle: 3 days",
        "Bug Rate: 15%",
        "Sprint Success: 60%",
        "Team Satisfaction: 6/10"
    ]
    
    for metric in metrics_before:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(16)  # Reduced from 18
        p.space_after = Pt(8)  # Reduced from 10
    
    # After metrics box
    right_left = Inches(6.9)
    right_content = slide.shapes.add_textbox(right_left, top, width, height)
    tf = right_content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "🎯 After 30 Days"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    metrics_after = [
        "PR Cycle: 4 hours ⬇️ 87%",
        "Bug Rate: 3% ⬇️ 80%",
        "Sprint Success: 92% ⬆️ 53%",
        "Team Satisfaction: 9/10 ⬆️ 50%"
    ]
    
    for metric in metrics_after:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    # Slide 10: Action Items
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🎬 Immediate Actions"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Do Today:"
    p.font.size = Pt(20)  # Reduced from 24
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    actions = [
        "☑️ Create your first project board",
        "☑️ Import standard label set",
        "☑️ Add issue templates",
        "☑️ Enable branch protection",
        "☑️ Configure PR template"
    ]
    
    for action in actions:
        p = tf.add_paragraph()
        p.text = action
        p.font.size = Pt(17)  # Reduced from 20
        p.space_after = Pt(6)  # Reduced from 8
    
    p = tf.add_paragraph()
    p.text = "\nResources:"
    p.font.size = Pt(19)  # Reduced from 22
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "📦 Templates: github.com/hanzlahabib/github-workflow-demo"
    p.font.size = Pt(15)  # Reduced from 18
    
    # Slide 11: Q&A
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💡 Questions & Discussion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Common Questions:"
    p.font.size = Pt(19)  # Reduced from 22
    p.font.bold = True
    
    questions = [
        "✅ How do we handle urgent hotfixes?",
        "✅ What about multiple repositories?",
        "✅ How to migrate existing issues?",
        "✅ What's the actual ROI?"
    ]
    
    for q in questions:
        p = tf.add_paragraph()
        p.text = q
        p.font.size = Pt(16)  # Reduced from 18
        p.space_after = Pt(8)  # Reduced from 10
    
    p = tf.add_paragraph()
    p.text = "\nContact & Support:"
    p.font.size = Pt(19)
    p.font.bold = True
    p.space_before = Pt(15)  # Reduced from 20
    
    contact_info = [
        "GitHub: @hanzlahabib | @amna",
        "Demo: github.com/hanzlahabib/github-workflow-demo",
        "Office Hours: Thursdays 2-3 PM"
    ]
    
    for info in contact_info:
        p = tf.add_paragraph()
        p.text = info
        p.font.size = Pt(15)  # Reduced from 18
    
    # Slide 12: Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🙏 Thank You!"
    title.text_frame.paragraphs[0].font.size = Pt(46)  # Reduced from 60
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Ready to Transform Your Workflow?\n\n🚀 Start Your 30-Day Journey Today\n\nLet's build better software, together."
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)  # Reduced from 24
    
    # Save the presentation
    prs.save('GitHub_Workflow_Best_Practices_Improved.pptx')
    print("✅ Improved PowerPoint presentation created successfully!")
    print("📁 File saved as: GitHub_Workflow_Best_Practices_Improved.pptx")
    print("🎯 Optimized for better visibility and content fitting")

if __name__ == "__main__":
    create_presentation()